from pptx import Presentation
import json
import os

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

ppt_path = os.path.join(
    BASE_DIR,
    "data",
    "Dataset_project_repository 1.pptx"
)

prs = Presentation(ppt_path)

slides_data = []

for i, slide in enumerate(prs.slides):

    slide_text = []

    for shape in slide.shapes:

        if hasattr(shape, "text"):

            text = shape.text.strip()

            if text:
                slide_text.append(text)

    slides_data.append({
        "slide_number": i + 1,
        "content": "\n".join(slide_text)
    })

output_dir = os.path.join(BASE_DIR, "output")
os.makedirs(output_dir, exist_ok=True)

output_path = os.path.join(output_dir, "slides.json")

with open(output_path, "w", encoding="utf-8") as f:
    json.dump(slides_data, f, indent=4, ensure_ascii=False)

print("✅ PowerPoint Parsed Successfully!")